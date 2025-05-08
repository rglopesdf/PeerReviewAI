# src/reference_parser.py

from pptx import Presentation
import os

class ReferenceParser:
    def __init__(self):
        pass

    def extract_text_from_pptx(self, pptx_path):
        """Extracts text from a given PPTX file.

        Args:
            pptx_path (str): The absolute path to the PPTX file.

        Returns:
            str: The extracted text content from the PPTX. Returns an empty string if extraction fails or file not found.
        """
        if not os.path.exists(pptx_path):
            print(f"Error: PPTX file not found at {pptx_path}")
            return ""
        
        try:
            prs = Presentation(pptx_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
                text_runs.append("\n\n--- Slide Break ---\n\n") # Add a separator for clarity between slides
            return " ".join(text_runs)
        except Exception as e:
            print(f"Error extracting text from PPTX {pptx_path}: {e}")
            return ""

if __name__ == '__main__':
    # Example Usage (for testing purposes)
    # This part would typically be run in an environment where a PPTX is available.
    
    # Create a dummy reference_materials directory for testing structure
    if not os.path.exists("/home/ubuntu/academic_evaluator/reference_materials"): 
        os.makedirs("/home/ubuntu/academic_evaluator/reference_materials")
    
    # A real PPTX would be needed here for actual testing.
    # dummy_pptx_path = "/home/ubuntu/academic_evaluator/reference_materials/sample.pptx"
    # # You would need to create a sample.pptx file for this to work
    # # For example, by saving an empty presentation as sample.pptx in the specified path.

    parser = ReferenceParser()
    print("ReferenceParser class defined. To test, place a PPTX in academic_evaluator/reference_materials/ and call extract_text_from_pptx.")
    # print(f"Please place a sample PPTX at {dummy_pptx_path} for testing.")
    # extracted_content = parser.extract_text_from_pptx(dummy_pptx_path)
    # if extracted_content:
    #     print(f"Extracted content from {dummy_pptx_path}:\n{extracted_content[:500]}...")
    # else:
    #     print(f"Could not extract content from {dummy_pptx_path}.")

